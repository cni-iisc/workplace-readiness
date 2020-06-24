# Copyright [2020] [Indian Institute of Science, Bangalore]
# SPDX-License-Identifier: Apache-2.0

import cv2
from datetime import datetime
import numpy as np
import os
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

def read_img_file(file_location):
    if os.path.isfile(file_location):
        im = cv2.imread(file_location)
        im = cv2.cvtColor(cv2.resize(im, dsize=(0,0), fx=0.3, fy=0.3), cv2.COLOR_RGB2RGBA)
    else:
        im = np.zeros((810, 1440, 4), dtype=np.uint8)
    return (im)        
         
def concat_tile(im_list_2d):
    return cv2.vconcat([cv2.hconcat(im_list_h) for im_list_h in im_list_2d])

def combine_image(plots_folder, file_position_map):
    if (os.path.isfile(plots_folder+file_position_map[str(1)])):
        im1 = read_img_file(plots_folder+file_position_map[str(1)]) 
        im2 = read_img_file(plots_folder+file_position_map[str(2)]) 
        im3 = read_img_file(plots_folder+file_position_map[str(3)]) 
        im4 = read_img_file(plots_folder+file_position_map[str(4)]) 
        im5 = read_img_file(plots_folder+file_position_map[str(5)]) 
        im6 = read_img_file(plots_folder+file_position_map[str(6)]) 
        im7 = read_img_file(plots_folder+file_position_map[str(7)]) 
        im8 = read_img_file(plots_folder+file_position_map[str(8)]) 
        im9 = read_img_file(plots_folder+file_position_map[str(9)]) 
        im10 = read_img_file(plots_folder+file_position_map[str(10)]) 

        img_height, img_width, channel = im10.shape
        n_channels = 4
        imt = np.zeros((img_height, img_width, n_channels), dtype=np.uint8)

        im_tile = concat_tile([[imt, imt, im10],
                               [im1, im2, im3],
                               [im4, im5, im6],
                               [im7, im8, im9]])
        cv2.imwrite(plots_folder+'opencv_concat_tile.png', im_tile)
    return (True) 

prs = Presentation('Daily_Reports/workplace_dailyReport_23_06.pptx')

file_position_map = {'1': 'Infrastructure.png', '2': 'Epidemic related: Precautions.png', '3': 'Epidemic related: Awareness and readiness.png',
                     '4': 'Epidemic related: Advertisement and outreach.png', '5': 'Transportation.png', '6': 'Employee interactions: Mobility.png',
                    '7': 'Employee interactions: Meetings.png', '8': 'Employee interactions: Outside contacts.png',
                    '9': 'Canteen_pantry.png', '10': 'Hygiene and sanitation.png', '11': 'Total.png'}

slide_to_pic = { '2': 'New_users', '3': 'Revisiting_users', '4': 'User_distibution_thus far_(Category_wise)', '5': 'User_distibution_today_(Category_wise)', '6': 'scores_hist/Total', '7': 'opencv_concat_tile', '8': 'scores_hist/opencv_concat_tile'}

for i in range(9):
    plots_folder = 'analytics/' + str(i+1) + '/'
    combine_image(plots_folder, file_position_map)

combine_image('analytics/scores_hist/', file_position_map)

im1 = read_img_file('analytics/1/Total.png') 
im2 = read_img_file('analytics/2/Total.png') 
im3 = read_img_file('analytics/3/Total.png') 
im4 = read_img_file('analytics/4/Total.png') 
im5 = read_img_file('analytics/5/Total.png') 
im6 = read_img_file('analytics/6/Total.png') 
im7 = read_img_file('analytics/7/Total.png') 
im8 = read_img_file('analytics/8/Total.png') 
im9 = read_img_file('analytics/9/Total.png') 

im_tile = concat_tile([[im1, im2, im3],
                       [im4, im5, im6],
                       [im7, im8, im9]])
cv2.imwrite('analytics/opencv_concat_tile.png', im_tile)

today = datetime.now()

slide_number = 0
for slide in prs.slides:
    slide_number += 1
    if (slide_number == 1):
        counter = 0
        for shape in slide.shapes:
            counter += 1
            if (shape.has_text_frame and counter==2):
                text_frame = shape.text_frame
                cur_text = text_frame.paragraphs[0].runs[0].text
                repl_str = 'Daily report: '+datetime.strftime(today, '%d %B %Y')
                new_text = cur_text.replace(str('Daily report: 23 June 2020'), str(repl_str))
                text_frame.paragraphs[0].runs[0].text = new_text 

    if (slide_number >= 2 and slide_number < 10):
        for shape in slide.shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                replace_img_path = 'analytics/'+slide_to_pic[str(slide_number)]+'.png'
                if (os.path.isfile(replace_img_path)):

                    imgPic = shape._pic
                    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                    imgPart = slide.part.related_parts[imgRID]

                    with open(replace_img_path, 'rb') as f:
                        rImgBlob = f.read()

                    #replace
                    imgPart._blob = rImgBlob
    
    if (slide_number >= 10):
        plots_folder = 'analytics/' + str(slide_number-9) + '/'
        for shape in slide.shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                replace_img_path = plots_folder+'opencv_concat_tile.png'
                if (os.path.isfile(replace_img_path)):

                    imgPic = shape._pic
                    imgRID = imgPic.xpath('./p:blipFill/a:blip/@r:embed')[0]
                    imgPart = slide.part.related_parts[imgRID]

                    with open(replace_img_path, 'rb') as f:
                        rImgBlob = f.read()

                    #replace
                    imgPart._blob = rImgBlob

prs.save('Daily_Reports/Workplace_Readiness_Daily_Report_'+datetime.strftime(today, '%d_%m')+'.pptx')
